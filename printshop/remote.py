# Flask 远程打印服务
from flask import Flask, request, jsonify, render_template, redirect, url_for, flash, session
from flask_login import LoginManager, login_user, logout_user, login_required, current_user, login_url
from flask_wtf.csrf import CSRFProtect, generate_csrf
from flask_wtf import FlaskForm
from wtforms import StringField, PasswordField
from wtforms.validators import DataRequired, Email
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from sqlalchemy import desc
import printer
import os
from models import db, User, RedemptionCode, PrintJob, OcrJob
import secrets
from datetime import datetime, timedelta
from flask_wtf.csrf import CSRFError
import pyotp
import qrcode
import io
import base64

app = Flask(__name__)
csrf = CSRFProtect(app)


@app.errorhandler(CSRFError)
def handle_csrf_error(e):
    # 返回友好的错误信息；如果是JSON请求则返回JSON
    if request.headers.get('Content-Type') == 'application/json':
        return jsonify({"error": "CSRF token missing or invalid"}), 400
    flash('CSRF 令牌缺失或无效，请重试', 'danger')
    return redirect(request.referrer or url_for('login'))

class LoginForm(FlaskForm):
    username = StringField('用户名', validators=[DataRequired()])
    password = PasswordField('密码', validators=[DataRequired()])

class Verify2FAForm(FlaskForm):
    otp = StringField('验证码', validators=[DataRequired()])
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///' + os.path.join(os.path.abspath(os.path.dirname(__file__)), 'printshop.db')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['SECRET_KEY'] = secrets.token_hex(16)

# 初始化数据库
db.init_app(app)
with app.app_context():
    db.create_all()
    # 创建默认管理员（如果不存在）
    try:
        admin = User.query.filter_by(username='root').first()
        if not admin:
            admin = User(username='root', password=generate_password_hash('123456'), balance=0.0, is_admin=True, has_unlimited=True)
            db.session.add(admin)
            db.session.commit()
            print('已创建默认管理员: root / 123456')
    except Exception as e:
        db.session.rollback()
        print(f'创建管理员失败: {e}')
    # 确保 OCR 结果目录存在
    os.makedirs(os.path.join('static', 'ocr_results'), exist_ok=True)
# OCR 配置
app.config['OCR_MODEL'] = 'minerU'  # 默认使用 minerU（需在后台配置可执行路径或API）
app.config['OCR_MODEL_PATH'] = None  # minerU 可执行文件或服务的路径/URL
app.config['OCR_PRICE_PER_PAGE'] = 1.0  # 每页 1 元

# 初始化登录管理
login_manager = LoginManager()
login_manager.init_app(app)

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

def is_json_request():
    return request.headers.get('Content-Type') == 'application/json'


def ocr_process_image(image_path):
    """对单张图片执行 OCR，优先使用 minerU（通过可执行路径或API），否则尝试 pytesseract（若安装）。返回识别到的文本。"""
    model = app.config.get('OCR_MODEL')
    model_path = app.config.get('OCR_MODEL_PATH')
    # minerU 调用（示例：假设 minerU 可执行文件接受文件路径并输出文本）
    if model == 'minerU' and model_path:
        try:
            import subprocess
            res = subprocess.run([model_path, image_path], capture_output=True, text=True, check=True)
            return res.stdout.strip()
        except Exception as e:
            print(f'minerU OCR 调用失败: {e}')

    # 回退到 pytesseract
    try:
        from PIL import Image
        import pytesseract
        text = pytesseract.image_to_string(Image.open(image_path), lang='chi_sim+eng')
        return text
    except Exception as e:
        print(f'pytesseract OCR 不可用或失败: {e}')

    raise RuntimeError('未配置可用的 OCR 引擎（minerU 或 pytesseract）')

# 用户认证路由
@app.route('/login', methods=['GET', 'POST'])
def login():
    form = LoginForm()
    if form.validate_on_submit():
        user = User.query.filter_by(username=form.username.data).first()
        if not user or not check_password_hash(user.password, form.password.data):
            flash('用户名或密码错误', 'danger')
            return redirect(url_for('login'))
        
        login_user(user)
        if user.is_2fa_enabled:
            session['user_id'] = user.id
            logout_user()
            return redirect(url_for('verify_2fa'))
        return redirect(url_for('dashboard'))
    
    if is_json_request() and request.method == 'POST':
        data = request.get_json()
        user = User.query.filter_by(username=data.get('username')).first()
        if not user or not check_password_hash(user.password, data.get('password')):
            return jsonify({"error": "Invalid username or password"}), 401
        login_user(user)
        return jsonify({"status": "success", "user_id": user.id})
    
    return render_template('login.html', form=form)

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        try:
            data = request.form if not is_json_request() else request.get_json()
            
            if not data.get('username') or not data.get('password'):
                if is_json_request():
                    return jsonify({"error": "Username and password required"}), 400
                flash('用户名和密码不能为空', 'danger')
                return redirect(url_for('register'))
            
            if User.query.filter_by(username=data['username']).first():
                if is_json_request():
                    return jsonify({"error": "Username already exists"}), 400
                flash('用户名已存在', 'danger')
                return redirect(url_for('register'))
            
            user = User(
                username=data['username'],
                password=generate_password_hash(data['password']),
                balance=0.0
            )
            db.session.add(user)
            db.session.commit()
            
            if is_json_request():
                return jsonify({"status": "success", "user_id": user.id})
            flash('注册成功，请登录', 'success')
            return redirect(url_for('login'))
        except Exception as e:
            db.session.rollback()
            if is_json_request():
                return jsonify({"error": str(e)}), 500
            flash(f'注册失败: {str(e)}', 'danger')
            return redirect(url_for('register'))
    
    return render_template('register.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    if is_json_request():
        return jsonify({"status": "success"})
    return redirect(url_for('login'))

# 用户功能路由
@app.route('/dashboard')
@login_required
def dashboard():
    recent_jobs = PrintJob.query.filter_by(user_id=current_user.id)\
        .order_by(desc(PrintJob.created_at))\
        .limit(5)\
        .all()
    return render_template('dashboard.html', recent_jobs=recent_jobs)

@app.route('/redeem', methods=['GET', 'POST'])
@login_required
def redeem():
    if request.method == 'POST':
        code = request.form.get('code') if not is_json_request() else request.get_json().get('code')
        redemption = RedemptionCode.query.filter_by(code=code, is_used=False).first()
        
        if not redemption:
            if is_json_request():
                return jsonify({"error": "Invalid or used code"}), 400
            flash('无效或已使用的兑换码', 'danger')
            return redirect(url_for('redeem'))
        
        # 如果兑换码为无限余额，赋予用户无限权限
        if getattr(redemption, 'is_unlimited', False):
            current_user.has_unlimited = True
        else:
            current_user.balance += redemption.amount
        redemption.is_used = True
        redemption.used_by = current_user.id
        redemption.used_at = datetime.utcnow()
        db.session.commit()
        
        if is_json_request():
            return jsonify({"status": "success", "new_balance": current_user.balance})
        flash(f'成功充值 ¥{redemption.amount:.2f}', 'success')
        return redirect(url_for('redeem'))
    
    return render_template('redeem.html')

@app.route('/history')
@login_required
def print_history():
    page = request.args.get('page', 1, type=int)
    search = request.args.get('search', '')
    start_date = request.args.get('start_date', '')
    end_date = request.args.get('end_date', '')

    query = PrintJob.query.filter_by(user_id=current_user.id)
    if search: query = query.filter(PrintJob.file_name.contains(search))
    if start_date: query = query.filter(PrintJob.created_at >= start_date)
    if end_date: query = query.filter(PrintJob.created_at <= end_date + ' 23:59:59')
    
    print_jobs = query.order_by(desc(PrintJob.created_at)).paginate(page=page, per_page=10)
    return render_template('history.html', print_jobs=print_jobs)

# 管理员API
@app.route('/admin/generate_code', methods=['POST'])
def generate_code():
    if request.headers.get('X-API-KEY') != API_KEY:
        return jsonify({"error": "Invalid API key"}), 401
    
    data = request.get_json()
    code = RedemptionCode(
        code=secrets.token_hex(8).upper(),
        amount=data.get('amount', 10.0)
    )
    db.session.add(code)
    db.session.commit()
    return jsonify({"code": code.code, "amount": code.amount})

# 打印功能
def generate_preview(file):
    """生成文件预览"""
    if not allowed_file(file.filename):
        return None
        
    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, f"preview_{filename}")
    file.save(filepath)
    
    try:
        # PDF预览
        if filename.lower().endswith('.pdf'):
            from PyPDF2 import PdfReader
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(filepath, first_page=1, last_page=1)
            except Exception:
                images = []
            preview_path = os.path.join(UPLOAD_FOLDER, f"preview_{os.path.splitext(filename)[0]}.jpg")
            if images:
                images[0].save(preview_path, 'JPEG')
            else:
                # 无法生成 PDF 预览，则使用空占位图或返回 None
                return None
            return preview_path
        
        # 图片预览 (直接使用原图)
        return filepath
    except Exception as e:
        print(f"生成预览失败: {str(e)}")
        return None
    finally:
        if os.path.exists(filepath) and filepath.endswith('.pdf'):
            os.remove(filepath)

@app.route('/preview', methods=['POST'])
@login_required
def handle_preview():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    preview_path = generate_preview(file)
    if not preview_path:
        return jsonify({"error": "Preview generation failed"}), 400
    
    # 返回预览图片URL
    preview_url = url_for('static', filename=f'uploads/{os.path.basename(preview_path)}', _external=True)
    return jsonify({"preview_url": preview_url})

@app.route('/print', methods=['POST'])
@login_required
def handle_print():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if not allowed_file(file.filename):
        return jsonify({"error": "File type not allowed"}), 400
        
    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    file.save(filepath)
    
    try:
        # 计算费用
        ext = os.path.splitext(filename)[1].lower()
        if ext == '.pdf':
            from PyPDF2 import PdfReader
            pages = len(PdfReader(filepath).pages)
            cost = pages * 0.5
        elif ext in ['.ppt', '.pptx', '.doc', '.docx']:
            # 尝试通过多种方式获取 Office 文档页/幻灯片数，若失败则按默认计价
            pages = get_office_page_count(filepath, ext)
            if pages is None:
                print("获取Office文档页数失败或未安装依赖，按默认1页计价")
                cost = 1.0
                pages = 1
            else:
                cost = pages * 0.5
        elif ext in ['.cpp', '.py', '.txt']:
            # 代码文件按50行=1页计算
            from printer import count_code_lines
            lines = count_code_lines(filepath)
            pages = max(1, lines // 50)  # 至少按1页计算
            cost = pages * 0.5
        else:
            cost = 1.0
        
        if current_user.balance < cost:
            raise ValueError("Insufficient balance")
        
        # 打印文件
        printer_name = request.form.get('printer_name', None)
        if not printer.print_file(filepath, printer_name):
            raise ValueError("Print failed")
        
        # 记录打印任务
        current_user.balance -= cost
        job = PrintJob(
            user_id=current_user.id,
            file_name=filename,
            pages=pages if ext == '.pdf' else 1,
            cost=cost,
            status='completed'
        )
        db.session.add(job)
        db.session.commit()
        
        return jsonify({
            "status": "success", 
            "cost": cost,
            "balance": current_user.balance
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        if os.path.exists(filepath):
            os.remove(filepath)

# 系统路由
@app.route('/setup_2fa', methods=['GET', 'POST'])
@login_required
def setup_2fa():
    if request.method == 'POST':
        otp = request.form.get('otp')
        if not otp or not current_user.otp_secret:
            flash('无效请求', 'danger')
            return redirect(url_for('setup_2fa'))
        
        totp = pyotp.TOTP(current_user.otp_secret)
        if not totp.verify(otp):
            flash('验证码错误', 'danger')
            return redirect(url_for('setup_2fa'))
        
        current_user.is_2fa_enabled = True
        db.session.commit()
        flash('2FA已成功启用', 'success')
        return redirect(url_for('dashboard'))
    
    if not current_user.otp_secret:
        current_user.otp_secret = pyotp.random_base32()
        db.session.commit()
    
    totp = pyotp.TOTP(current_user.otp_secret)
    uri = totp.provisioning_uri(name=current_user.username, issuer_name="打印店系统")
    
    # 生成QR码
    img = qrcode.make(uri)
    buffered = io.BytesIO()
    img.save(buffered)
    qr_code_url = "data:image/png;base64," + base64.b64encode(buffered.getvalue()).decode()
    
    return render_template('setup_2fa.html', 
                         qr_code_url=qr_code_url,
                         secret=current_user.otp_secret)

@app.route('/verify_2fa', methods=['GET', 'POST'])
def verify_2fa():
    if 'user_id' not in session:
        return redirect(url_for('login'))
    
    user = User.query.get(session['user_id'])
    if not user:
        return redirect(url_for('login'))
    
    form = Verify2FAForm()
    if form.validate_on_submit():
        totp = pyotp.TOTP(user.otp_secret)
        if totp.verify(form.otp.data):
            session.pop('user_id', None)
            login_user(user)
            return redirect(url_for('dashboard'))
        flash('验证码错误', 'danger')
    
    return render_template('verify_2fa.html', form=form)

class ForgotPasswordForm(FlaskForm):
    identifier = StringField('用户名或注册邮箱', validators=[DataRequired()])

class ResetPasswordForm(FlaskForm):
    new_password = PasswordField('新密码', validators=[DataRequired()])
    confirm_password = PasswordField('确认新密码', validators=[DataRequired()])

class ChangePasswordForm(FlaskForm):
    current_password = PasswordField('当前密码', validators=[DataRequired()])
    new_password = PasswordField('新密码', validators=[DataRequired()])
    confirm_password = PasswordField('确认新密码', validators=[DataRequired()])

@app.route('/change_password', methods=['GET', 'POST'])
@login_required
def change_password():
    form = ChangePasswordForm()
    if form.validate_on_submit():
        if not check_password_hash(current_user.password, form.current_password.data):
            flash('当前密码错误', 'danger')
            return redirect(url_for('change_password'))
        
        if form.new_password.data != form.confirm_password.data:
            flash('新密码不匹配', 'danger')
            return redirect(url_for('change_password'))
            
        current_user.password = generate_password_hash(form.new_password.data)
        db.session.commit()
        flash('密码已成功修改', 'success')
        return redirect(url_for('dashboard'))
    
    return render_template('change_password.html', form=form)

@app.route('/disable_2fa', methods=['POST'])
@login_required
def disable_2fa():
    # 强化：要求当前密码 + OTP
    current_pwd = request.form.get('current_password')
    otp = request.form.get('otp_disable')
    if not current_pwd or not otp:
        flash('请提供当前密码和 OTP', 'danger')
        return redirect(url_for('setup_2fa'))

    if not check_password_hash(current_user.password, current_pwd):
        flash('当前密码错误', 'danger')
        return redirect(url_for('setup_2fa'))

    if not current_user.otp_secret:
        flash('未配置 OTP', 'danger')
        return redirect(url_for('setup_2fa'))

    totp = pyotp.TOTP(current_user.otp_secret)
    if not totp.verify(otp):
        flash('OTP 验证失败', 'danger')
        return redirect(url_for('setup_2fa'))

    current_user.otp_secret = None
    current_user.is_2fa_enabled = False
    db.session.commit()
    flash('2FA已禁用', 'success')
    return redirect(url_for('dashboard'))

@app.route('/')
def index():
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))
    return redirect(url_for('login'))

@app.route('/forgot_password', methods=['GET', 'POST'])
def forgot_password():
    form = ForgotPasswordForm()
    if form.validate_on_submit():
        # 支持通过用户名或邮箱（目前模型仅有 username 字段）进行重置
        identifier = form.identifier.data
        user = User.query.filter_by(username=identifier).first()
        if user:
            # 生成重置令牌(示例，实际应使用更安全的方法)
            token = secrets.token_urlsafe(32)
            user.reset_token = token
            user.reset_token_expires = datetime.utcnow() + timedelta(hours=1)
            db.session.commit()
            
            # 这里应该发送邮件，示例中仅打印到控制台
            reset_url = url_for('reset_password', token=token, _external=True)
            print(f"密码重置链接: {reset_url}")  # 实际应用中应发送邮件
            flash('如果该邮箱已注册，已发送重置链接', 'success')
        else:
            # 防止用户枚举攻击，无论用户是否存在都显示相同消息
            flash('如果该邮箱已注册，已发送重置链接', 'success')
        return redirect(url_for('login'))
    
    return render_template('forgot_password.html', form=form)

@app.route('/reset_password/<token>', methods=['GET', 'POST'])
def reset_password(token):
    user = User.query.filter_by(reset_token=token).first()
    
    if not user or user.reset_token_expires < datetime.utcnow():
        flash('无效或过期的重置链接', 'danger')
        return redirect(url_for('login'))
    
    form = ResetPasswordForm()
    if form.validate_on_submit():
        user.password = generate_password_hash(form.new_password.data)
        user.reset_token = None
        user.reset_token_expires = None
        db.session.commit()
        flash('密码已重置，请使用新密码登录', 'success')
        return redirect(url_for('login'))
    
    return render_template('reset_password.html', form=form, token=token)

@app.route('/health')
def health_check():
    return jsonify({"status": "running", "service": "printshop"})

# 配置
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx', 'doc', 'docx', 'png', 'jpg', 'jpeg', 'bmp', 'gif', 'cpp', 'py', 'txt'}
API_KEY = "printshop123"

# 确保上传目录存在
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def get_office_page_count(filepath, ext):
    """尝试多种方式获取 Office 文档的页/幻灯片数：
    1. 首选 Windows COM (pywin32)；
    2. 其次尝试 python-pptx / python-docx（若已安装）；
    3. 最后返回 None，由调用方决定默认计价。
    """
    # 1) 使用 win32com（仅限 Windows，需安装 pywin32）
    try:
        from win32com.client import Dispatch
        app = None
        try:
            if ext in ['.ppt', '.pptx']:
                app = Dispatch('PowerPoint.Application')
                pres = app.Presentations.Open(os.path.abspath(filepath))
                pages = pres.Slides.Count
                pres.Close()
            else:
                app = Dispatch('Word.Application')
                doc = app.Documents.Open(os.path.abspath(filepath))
                pages = doc.ComputeStatistics(2)  # wdStatisticPages
                doc.Close()
            return int(pages)
        finally:
            if app:
                try:
                    app.Quit()
                except Exception:
                    pass
    except Exception as e:
        print(f"win32com unavailable or failed: {e}")

    # 2) 尝试 python-pptx / python-docx
    try:
        if ext in ['.ppt', '.pptx']:
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                return len(prs.slides)
            except Exception as e:
                print(f"python-pptx unavailable or failed: {e}")
                return None
        else:
            # python-docx 无法可靠地获取页数，故返回 None
            return None
    except Exception as e:
        print(f"office page count fallback failed: {e}")

    return None

def run_server(host='0.0.0.0', port=5000):
    app.run(host=host, port=port)

if __name__ == '__main__':
    run_server()


@app.route('/admin/ocr_config', methods=['GET', 'POST'])
@login_required
def admin_ocr_config():
    if not current_user.is_admin:
        flash('需要管理员权限', 'danger')
        return redirect(url_for('dashboard'))

    if request.method == 'POST':
        model = request.form.get('model')
        model_path = request.form.get('model_path')
        app.config['OCR_MODEL'] = model or app.config.get('OCR_MODEL')
        app.config['OCR_MODEL_PATH'] = model_path or app.config.get('OCR_MODEL_PATH')
        flash('OCR 配置已更新', 'success')
        return redirect(url_for('admin_ocr_config'))

    return render_template('admin_ocr_config.html', model=app.config.get('OCR_MODEL'), model_path=app.config.get('OCR_MODEL_PATH'))


@app.route('/admin/codes', methods=['GET', 'POST'])
@login_required
def admin_codes():
    if not current_user.is_admin:
        flash('需要管理员权限', 'danger')
        return redirect(url_for('dashboard'))

    if request.method == 'POST':
        code = request.form.get('code') or secrets.token_hex(8).upper()
        amount = request.form.get('amount')
        is_unlimited = True if request.form.get('is_unlimited') == 'on' else False
        try:
            amt = float(amount) if amount and not is_unlimited else 0.0
        except Exception:
            amt = 0.0

        new_code = RedemptionCode(code=code, amount=amt, is_unlimited=is_unlimited)
        db.session.add(new_code)
        db.session.commit()
        flash(f'已创建兑换码 {new_code.code}', 'success')
        return redirect(url_for('admin_codes'))

    codes = RedemptionCode.query.order_by(desc(RedemptionCode.created_at)).limit(50).all()
    return render_template('admin_codes.html', codes=codes)


@app.route('/ocr', methods=['POST'])
@login_required
def handle_ocr():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if not allowed_file(file.filename):
        return jsonify({'error': 'File type not allowed'}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    file.save(filepath)
    try:
        ext = os.path.splitext(filename)[1].lower()
        pages = 1
        text_output = ''
        # 图片直接 OCR
        if ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif']:
            pages = 1
            text_output = ocr_process_image(filepath)
        elif ext == '.pdf':
            # 尝试将 PDF 每页转为图片并 OCR（需要 pdf2image + poppler）
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(filepath)
                pages = len(images)
                texts = []
                for i, img in enumerate(images):
                    tmp = os.path.join(UPLOAD_FOLDER, f'_ocr_page_{i}.png')
                    img.save(tmp, 'PNG')
                    texts.append(ocr_process_image(tmp))
                    os.remove(tmp)
                text_output = '\n'.join(texts)
            except Exception as e:
                return jsonify({'error': f'PDF OCR 失败: {e}'}), 500
        elif ext in ['.ppt', '.pptx']:
            # 尝试提取 ppt 文本（使用 python-pptx），若不可用则返回错误
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                pages = len(prs.slides)
                texts = []
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            slide_text.append(shape.text)
                    texts.append('\n'.join(slide_text))
                text_output = '\n\n'.join(texts)
            except Exception as e:
                return jsonify({'error': f'PPTX 文本提取失败: {e}'}), 500
        else:
            return jsonify({'error': '不支持的文件类型进行 OCR'}), 400

        # 如果用户是无限余额用户，不扣费
        price_per = app.config.get('OCR_PRICE_PER_PAGE', 1.0)
        # 管理员或已被赋予无限权限的用户不扣费
        is_unlimited_user = getattr(current_user, 'has_unlimited', False) or getattr(current_user, 'is_admin', False)
        cost = 0.0 if is_unlimited_user else pages * price_per
        if not is_unlimited_user and current_user.balance < cost:
            return jsonify({'error': 'Insufficient balance'}), 402

        # 扣费并保存任务，同时把 OCR 结果写入静态文件以便下载
        current_user.balance -= cost
        # 生成唯一文件名并保存文本结果
        result_fname = f'ocr_{secrets.token_hex(8)}.txt'
        result_path = os.path.join('static', 'ocr_results', result_fname)
        with open(result_path, 'w', encoding='utf-8') as f:
            f.write(text_output)

        job = OcrJob(user_id=current_user.id, file_name=filename, pages=pages, cost=cost, result_text=text_output, result_file=result_fname, status='completed')
        db.session.add(job)
        db.session.commit()

        download_url = url_for('static', filename=f'ocr_results/{result_fname}', _external=True)
        return jsonify({'status': 'success', 'pages': pages, 'cost': cost, 'text': text_output, 'download_url': download_url})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500
    finally:
        if os.path.exists(filepath):
            os.remove(filepath)
